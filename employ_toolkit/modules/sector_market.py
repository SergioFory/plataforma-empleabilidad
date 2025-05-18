# employ_toolkit/modules/sector_market.py
from datetime import date
import uuid
from pathlib import Path
import re

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

OUTPUT_DIR = Path("workspace")
OUTPUT_DIR.mkdir(exist_ok=True)


# --------------------------------------------------------------------------- #
# Helpers                                                                     #
# --------------------------------------------------------------------------- #
def _add_bullets(box, text):
    """Convierte líneas de texto en viñetas."""
    for line in text.splitlines():
        if not line.strip():
            continue
        p = box.add_paragraph()
        p.text = line.strip()
        p.level = 0
        p.font.size = Pt(14)


def _parse_salary_lines(text: str):
    """
    Parsea líneas como 'Senior 40-55k' → (labels, values)
    Toma el punto medio del rango. Si hay un solo número, lo usa tal cual.
    """
    labels, values = [], []
    for line in text.splitlines():
        if not line.strip():
            continue
        parts = line.split()
        label = parts[0]
        nums = re.findall(r"\d+(?:\.\d+)?", line)
        if not nums:
            continue
        val = float(nums[0]) if len(nums) == 1 else (float(nums[0]) + float(nums[1])) / 2
        labels.append(label)
        values.append(val)
    return labels, values


# --------------------------------------------------------------------------- #
# Main generator                                                              #
# --------------------------------------------------------------------------- #
def generate_sector_ppt(client, answers: dict) -> Path:
    """
    Genera una presentación PPTX con la información de sector / mercado
    y un gráfico de barras con los rangos salariales.
    Returns: Path al archivo generado.
    """
    prs = Presentation()

    # ---------------- Slide Título ---------------------------------------
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f"Sector & Mercado – {client.full_name}"
    title_slide.placeholders[1].text = (
        f"{answers['sector']} · {answers['region']}\n{date.today()}"
    )

    # ---------------- Slide Empresas / Roles -----------------------------
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "Empresas y Roles demandados"
    body1 = s1.shapes.placeholders[1].text_frame
    body1.clear()
    body1.text = "Empresas clave"
    _add_bullets(body1, answers["empresas"])
    p_roles = body1.add_paragraph(); p_roles.text = "\nRoles más demandados"; p_roles.font.bold = True
    _add_bullets(body1, answers["roles"])

    # ---------------- Slide Habilidades ----------------------------------
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Habilidades & Salarios"
    body2 = s2.shapes.placeholders[1].text_frame
    body2.clear()
    body2.text = "Habilidades técnicas top"
    _add_bullets(body2, answers["hards"])
    p_soft = body2.add_paragraph(); p_soft.text = "\nSoft skills top"; p_soft.font.bold = True
    _add_bullets(body2, answers["softs"])

    # ---------------- Slide Tendencias / Retos ---------------------------
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Tendencias y Retos"
    body3 = s3.shapes.placeholders[1].text_frame
    body3.clear()
    body3.text = "Tendencias"
    _add_bullets(body3, answers["tendencias"])
    p_ret = body3.add_paragraph(); p_ret.text = "\nRetos / pain-points"; p_ret.font.bold = True
    _add_bullets(body3, answers["retos"])

    # ---------------- Slide Gráfico Salarial -----------------------------
    labels, values = _parse_salary_lines(answers["salarios"])
    if labels:  # Solo crea la diapositiva si hay datos numéricos
        s4 = prs.slides.add_slide(prs.slide_layouts[5])  # título y contenido
        s4.shapes.title.text = "Comparativa de rangos salariales"

        chart_data = CategoryChartData()
        chart_data.categories = labels
        chart_data.add_series("Salario medio", values)

        x, y, cx, cy = Inches(1), Inches(2), Inches(8), Inches(4)
        chart = s4.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
        ).chart
        chart.category_axis.tick_labels.font.size = Pt(12)
        chart.value_axis.tick_labels.font.size = Pt(12)
        chart.value_axis.has_major_gridlines = False

    # ---------------- Guardar --------------------------------------------
    filename = OUTPUT_DIR / f"{client.full_name}_{uuid.uuid4().hex[:6]}_sector.pptx"
    prs.save(filename)
    return filename
