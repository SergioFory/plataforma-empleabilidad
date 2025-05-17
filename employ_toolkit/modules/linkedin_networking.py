from pathlib import Path
from pptx import Presentation

OUTPUT_DIR = Path("workspace")
OUTPUT_DIR.mkdir(exist_ok=True)

def networking_ppt(context):
    profile_name = context["intake"].full_name
    prs = Presentation()
    
    # Portada
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Networking en LinkedIn"
    slide.placeholders[1].text = f"Guía para {profile_name}"

    # Diapositiva perfil
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Optimiza tu Perfil"
    s.shapes.placeholders[1].text = "- Foto profesional\n- Titular con PROBLEMA + RESULTADO\n- About orientado a tu oferta de valor"

    # Diapositiva búsqueda
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Búsquedas Avanzadas"
    s.shapes.placeholders[1].text = "- Filtros por industria y cargo\n- Guardar búsquedas\n- Crear alertas"

    # Diapositiva mensajes
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Mensajes de Contacto en Frío"
    s.shapes.placeholders[1].text = "Ejemplo:\nHola {{nombre}}, vi que lideras {{equipo}} en {{empresa}} …"

    file_path = OUTPUT_DIR / f"networking_linkedin_{profile_name}.pptx"
    prs.save(file_path)
    print(f"✓ Presentación LinkedIn guardada en {file_path}")
    return file_path
